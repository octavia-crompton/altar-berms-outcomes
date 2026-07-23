"""Build an AGU poster (session H138) using the CUAHSI biennial poster as a
visual template: 48x36 in, four columns, Montserrat, teal/navy palette.

Content is the 775-berm remote-sensing study; the design system (canvas size,
fonts, colors, column grid, rounded panels) is inherited from
'CUAHSI biennial.pptx'.

Usage:  python scripts/build_agu_poster.py
Output: presentations/agu_poster_H138.pptx  (gitignored)
"""
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

REPO = Path(__file__).resolve().parents[1]
FIGS = REPO / "draft" / "overleaf" / "Figures"
OUTC = REPO / "figures" / "outcomes"
OUT = REPO / "presentations" / "agu_poster_H138.pptx"
TEMPLATE = REPO / "presentations" / "CUAHSI biennial.pptx"   # inherit its embedded Montserrat + 48x36 canvas

# ── Template design tokens ───────────────────────────────────────────────────
FONT = "Montserrat"
TEAL = RGBColor(0x14, 0x82, 0xA5)   # title / figure headers
NAVY = RGBColor(0x23, 0x50, 0x78)   # section headers
GREEN = RGBColor(0x5F, 0x7A, 0x66)  # takeaway accent
INK = RGBColor(0x22, 0x22, 0x22)    # body
MUTE = RGBColor(0x60, 0x70, 0x7A)   # captions
PANEL = RGBColor(0xF3, 0xF7, 0xF9)  # light panel fill
PANEL_LN = RGBColor(0xD3, 0xE0, 0xE6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

def IN(v): return Inches(v)


def set_theme_font(prs, name):
    """Set the deck's theme major/minor Latin fonts, so any text that falls back
    to the theme (instead of an explicit run font) still renders in *name*.
    The CUAHSI template overrides runs with Montserrat but leaves the theme as
    Times New Roman; matching the theme too makes the look robust."""
    from lxml import etree
    ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    for master in prs.slide_masters:
        for rel in master.part.rels.values():
            if "theme" not in rel.reltype:
                continue
            part = rel.target_part
            root = etree.fromstring(part.blob)
            for tag in ("majorFont", "minorFont"):
                for latin in root.findall(f".//a:{tag}/a:latin", ns):
                    latin.set("typeface", name)
            part._blob = etree.tostring(root, xml_declaration=True,
                                        encoding="UTF-8", standalone=True)

# 48 x 36 canvas, four 11-in columns
COL_X = [0.6, 12.55, 24.5, 36.45]
COL_W = 11.0
PAD = 0.4                      # inner panel padding
CONTENT_TOP = 7.3
CONTENT_BOT = 35.4
GAP = 0.45                     # vertical gap between panels


def text(slide, s, x, y, w, h, size, color, bold=False, align=PP_ALIGN.LEFT,
         italic=False, anchor=MSO_ANCHOR.TOP, line_spacing=1.05):
    tb = slide.shapes.add_textbox(IN(x), IN(y), IN(w), IN(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = s.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run(); r.text = ln
        r.font.name = FONT; r.font.size = Pt(size); r.font.bold = bold
        r.font.italic = italic; r.font.color.rgb = color
    return tb


def bullets(slide, items, x, y, w, h, size, color):
    tb = slide.shapes.add_textbox(IN(x), IN(y), IN(w), IN(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.05
        p.space_after = Pt(10)
        r = p.add_run(); r.text = "▪  " + it
        r.font.name = FONT; r.font.size = Pt(size); r.font.color.rgb = color
    return tb


def panel(slide, x, y, w, h):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, IN(x), IN(y), IN(w), IN(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = PANEL
    sh.line.color.rgb = PANEL_LN; sh.line.width = Pt(1.5)
    sh.shadow.inherit = False
    try:
        sh.adjustments[0] = 0.03
    except Exception:
        pass
    return sh


CPI = 5.1   # approx characters per inch of text width at body size (for wrapping)


class Column:
    """Collects section specs, then lays them out vertically justified."""
    def __init__(self, slide, xi):
        self.slide = slide
        self.x = COL_X[xi]
        self.specs = []

    def section(self, header=None, **spec):
        if header is not None:
            spec["header"] = header
        self.specs.append(spec)

    # -- measure one section: return (parts, total_height) -------------------
    def _layout(self, sp):
        inner_w = COL_W - 2 * PAD
        body_size = sp.get("body_size", 21)
        cy = PAD
        parts = [("header", cy, (sp["header"], sp.get("header_color", NAVY)))]
        cy += 1.15
        if sp.get("body"):
            b = sp["body"]
            nlines = sum(max(1, int(len(seg) / (inner_w * CPI)) + 1)
                         for seg in b.split("\n"))
            bh = nlines * (body_size / 72.0 * 1.32) + 0.12
            parts.append(("body", cy, (b, body_size, bh)))
            cy += bh + 0.12
        if sp.get("bullets_items"):
            items = sp["bullets_items"]
            nlines = sum(max(1, int(len(b) / (inner_w * CPI)) + 1) for b in items)
            bh = nlines * (body_size / 72.0 * 1.32) + 0.16 * len(items)
            parts.append(("bullets", cy, (items, body_size, bh)))
            cy += bh + 0.1
        if sp.get("image"):
            iw, ih = Image.open(sp["image"]).size
            disp_w = inner_w * sp.get("img_scale", 1.0)
            disp_h = disp_w * ih / iw
            parts.append(("image", cy, (sp["image"], disp_w, disp_h)))
            cy += disp_h + 0.14
        if sp.get("caption"):
            c = sp["caption"]
            ch = (int(len(c) / (inner_w * (CPI + 1))) + 1) * 0.30 + 0.12
            parts.append(("caption", cy, (c, ch)))
            cy += ch + 0.04
        if sp.get("takeaway"):
            t = sp["takeaway"]
            th = (int(len(t) / (inner_w * (CPI - 1.4))) + 1) * 0.42 + 0.12
            parts.append(("takeaway", cy, (t, th)))
            cy += th + 0.04
        cy += PAD - 0.15
        return parts, cy

    def _draw(self, sp, y, parts, total_h):
        x, w = self.x, COL_W
        inner_x, inner_w = x + PAD, COL_W - 2 * PAD
        panel(self.slide, x, y, w, total_h)
        for kind, ry, payload in parts:
            yy = y + ry
            if kind == "header":
                s, col = payload
                text(self.slide, s, inner_x, yy, inner_w, 1.0, 32, col, bold=False)
            elif kind == "body":
                s, sz, bh = payload
                text(self.slide, s, inner_x, yy, inner_w, bh, sz, INK)
            elif kind == "bullets":
                items, sz, bh = payload
                bullets(self.slide, items, inner_x, yy, inner_w, bh, sz, INK)
            elif kind == "image":
                img, dw, dh = payload
                self.slide.shapes.add_picture(str(img), IN(x + (w - dw) / 2),
                                              IN(yy), IN(dw), IN(dh))
            elif kind == "caption":
                s, ch = payload
                text(self.slide, s, inner_x, yy, inner_w, ch, 18, MUTE, italic=True)
            elif kind == "takeaway":
                s, th = payload
                text(self.slide, s, inner_x, yy, inner_w, th, 22, GREEN, bold=False)

    def render(self, verbose=False):
        laid = [self._layout(sp) for sp in self.specs]
        heights = [h for _, h in laid]
        n = len(heights)
        avail = CONTENT_BOT - CONTENT_TOP
        # gap = the slack distributed exactly across the n-1 inter-panel gaps,
        # clamped so panels neither collide nor drift too far apart.
        if n > 1:
            gap = (avail - sum(heights)) / (n - 1)
            gap = max(0.30, min(3.1, gap))
        else:
            gap = GAP
        y = CONTENT_TOP
        for sp, (parts, h) in zip(self.specs, laid):
            self._draw(sp, y, parts, h)
            y += h + gap
        bottom = y - gap
        if verbose:
            print(f"  col x={self.x}: sum_h={sum(heights):.1f} gap={gap:.2f} "
                  f"bottom={bottom:.1f} {'OVERFLOW' if bottom > CONTENT_BOT + 0.05 else 'ok'}")
        return bottom


def build():
    # Start from the CUAHSI template so the poster inherits its embedded
    # Montserrat, theme, 48x36 canvas, and background; then clear its shapes.
    prs = Presentation(str(TEMPLATE))
    s = prs.slides[0]
    for shp in list(s.shapes):
        shp._element.getparent().remove(shp._element)
    set_theme_font(prs, FONT)

    # ── Title band ───────────────────────────────────────────────────────────
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, IN(0), IN(0), IN(48), IN(6.6))
    band.fill.solid(); band.fill.fore_color.rgb = WHITE
    band.line.fill.background(); band.shadow.inherit = False
    rule = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, IN(0), IN(6.5), IN(48), IN(0.12))
    rule.fill.solid(); rule.fill.fore_color.rgb = TEAL
    rule.line.fill.background(); rule.shadow.inherit = False

    text(s, "Controls on berm structural condition and vegetation response "
            "in the Altar Valley, southern Arizona",
         1.0, 0.55, 46.0, 2.8, 74, TEAL, bold=False, align=PP_ALIGN.CENTER)
    text(s, "O. V. Crompton¹, A. N. Koop², S. Anderson³, S. Chaulagain³, "
            "D. A. Lapides⁴, M. B. Meles⁵, M. H. Nichols⁴",
         1.0, 4.15, 46.0, 0.9, 38, TEAL, bold=False, align=PP_ALIGN.CENTER)
    text(s, "¹USDA-ARS Hydrology & Remote Sensing Lab  ·  ²ORISE / USDA-ARS SCINet  ·  "
            "³Univ. of Arizona  ·  ⁴USDA-ARS Southwest Watershed Research Center  ·  "
            "⁵USDA-ARS Sustainable Ag. Water Systems      |      AGU H138 — Human Impacts on the Water Cycle",
         1.0, 5.05, 46.0, 0.8, 19, MUTE, align=PP_ALIGN.CENTER)

    # ── Column 1 — background, study area, open question ─────────────────────
    c1 = Column(s, 0)
    c1.section(
        "Legacy conservation earthworks",
        body="Beginning in the 1930s, the US Soil Conservation Service (now the "
             "NRCS) partnered with private landowners to fund and construct water "
             "management structures, resulting in thousands of earthen berms across "
             "southwestern rangelands. Water spreader berms were typically "
             "constructed as large earthen structures (often more than 100 m long "
             "and 3–5 m wide) along landscape contours, singly or in series. "
             "Built perpendicular to the flow direction, they reduce flow "
             "velocities, allowing time for runoff to infiltrate and increase soil "
             "moisture, which is critical for supporting vegetation.",
        body_size=21)
    c1.section(
        header="Study area: Altar Valley, southern Arizona",
        header_color=TEAL,
        body="A semi-arid, 247,000-ha watershed where >1,500 structures have been "
             "built since the early 1900s; over half of the 415 mm annual rainfall "
             "falls in the July–September monsoon. We analyze 775 mapped water "
             "spreader berms, classified intact, breached, or flanked (Google Earth "
             "+ LiDAR).",
        body_size=20,
        image=FIGS / "LocMap.png",
        caption="Berm locations colored by structural condition (green, intact; red, "
                "degraded) over shaded relief, with the stream network (blue) and "
                "study-area boundary (black). Insets: location in the SW US; aerial "
                "imagery of a representative berm cluster.")
    c1.section(
        "The open question", header_color=TEAL,
        body="Earthen berms fail over time: many are breached (runoff cuts through) "
             "or flanked (runoff routes around the end). Yet "
             "few studies have addressed how landscape controls — geomorphic "
             "position, soils, and local topography — influence berm condition and "
             "vegetation response. Likewise, whether berm failure affects vegetation "
             "response is unknown.",
        body_size=21)

    # ── Column 2 — vegetation-response metric and approach ───────────────────
    c2 = Column(s, 1)
    c2.section(
        "Quantifying vegetation response",
        header_color=TEAL,
        image=FIGS / "fig2_veg_response.png",
        caption="Vegetation response metric for an example berm.",
        takeaway="Following Crompton et al. (2025), vegetation response (ΔS) is the "
                 "percent difference in SAVI between areas 15–60 m upslope and "
                 "downslope of a berm, normalized by background conditions (areas "
                 ">200 m from any berm).")
    c2.section(
        "Approach",
        body="We use controlled logistic regression to isolate the independent "
             "contributions of landform, soil texture, soil development, berm "
             "length, surface slope, and flow accumulation to berm structural "
             "condition (intact vs. degraded) and vegetation response. Drop-one "
             "likelihood-ratio tests (ΔAIC > 2 = meaningful) flag each predictor's "
             "independent contribution, complemented by a Random Forest classifier "
             "(5-fold CV). Predictors derive from NRCS SSURGO soils and "
             "LiDAR-derived terrain; nearly half of berms (47%) exceed the 7% "
             "SAVI-difference threshold for a measurable response.",
        body_size=21)
    c2.section(
        "Predictive performance",
        header_color=TEAL,
        image=FIGS / "fig5_model_performance_importance.png",
        caption="Random Forest permutation importance for each outcome.",
        takeaway="Predictive accuracy is modest (CV AUC = 0.71 ± 0.03 for condition, "
                 "0.64 ± 0.05 for vegetation response), identifying broad "
                 "landscape-scale controls rather than mechanistic predictors.")

    # ── Column 3 — controls on each outcome ──────────────────────────────────
    c3 = Column(s, 2)
    c3.section(
        "Structural condition",
        header_color=TEAL,
        body="Of 775 berms, 458 (59%) are intact and 317 (41%) degraded — 203 "
             "flanked and 114 breached. Controlling for all predictors, condition "
             "is driven by berm length (ΔAIC +30.6, p<0.001), flow accumulation "
             "(+6.7, p=0.003), soil texture (+3.0, p=0.017), and B-horizon presence "
             "(+2.3, p=0.038); slope and landform are not independently significant.",
        body_size=20,
        image=FIGS / "fig6_berm_condition.png",
        caption="Six predictors of berm structural condition (intact vs. degraded).",
        takeaway="Berms are more likely to remain intact when they are shorter, "
                 "located in positions of lower flow accumulation, and built on "
                 "finer-textured soils.")
    c3.section(
        "Vegetation response",
        header_color=TEAL,
        body="Of 775 berms, 368 (47%) show a vegetation response and 407 (53%) do "
             "not. After controlling, only slope (ΔAIC +18.8, p<0.001) and soil "
             "texture (+5.6, p=0.006) remain significant. Soil texture predicts both "
             "outcomes in opposite directions — finer soils favor intact berms, "
             "coarser (sandy loam) soils favor vegetation response (the inverse "
             "texture effect).",
        body_size=20,
        image=FIGS / "fig7_vegetation_response.png",
        caption="Six predictors of vegetation response.",
        takeaway="Vegetation responses are larger on steeper slopes and "
                 "coarser-textured (sandy loam) soils.")

    # ── Column 4 — independence, confounding, conclusions ────────────────────
    c4 = Column(s, 3)
    c4.section(
        "Condition and response are independent",
        header_color=NAVY,
        image=FIGS / "fig8_veg_response_by_condition_texture.png",
        caption="Co-occurrence of structural condition and vegetation response.",
        takeaway="Structural condition and vegetation response are statistically "
                 "independent (φ = −0.02, p = 0.612): condition does not predict "
                 "response. Only sandy-loam berms differ (52% vs. 38%, p = 0.034).")
    c4.section(
        "Landform is confounded",
        header_color=TEAL,
        image=FIGS / "fig3_controlled_predictors.png",
        takeaway="Landform is univariately associated with vegetation response "
                 "(p<0.001) but drops to p=0.310 after controlling for slope and "
                 "soil texture — evidence of confounding.")
    c4.section(
        "Conclusions", header_color=NAVY,
        bullets_items=[
            "Structural condition and vegetation response are shaped by largely separate landscape controls.",
            "Condition is determined by berm length, flow accumulation, and soil texture; vegetation response is more sensitive to slope and soil texture.",
            "The two outcomes are statistically independent, so physical integrity and ecological function must be evaluated separately.",
            "These controls can be read from existing terrain and soil-survey data — a basis for prioritizing maintenance or removal across large berm inventories.",
        ], body_size=19)
    c4.section(
        "Acknowledgements & references", header_color=TEAL,
        body="Supported by USDA-NRCS CEAP–Grazing Lands (NR213A750023C013). Soil "
             "data: USDA-NRCS SSURGO. LiDAR: Pima County RFCD. Berm inventory: "
             "Nichols and Degginger (2021), Catena. Vegetation metric follows "
             "Crompton et al. (2025).",
        body_size=16)

    for col in (c1, c2, c3, c4):
        col.render(verbose=True)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    return prs


def main():
    prs = build()
    print(f"Saved {OUT}  ({len(prs.slides._sldIdLst)} slide, 48x36 in)")


if __name__ == "__main__":
    main()
