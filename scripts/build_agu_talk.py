"""Build a minimal AGU talk deck (session H138 — Human Impacts on the Water Cycle).

Minimal design: one idea per slide, the takeaway *is* the title, the figure does
the talking. Optional one-line subtitle; no bullet walls.

Usage:
    python scripts/build_agu_talk.py
Output:
    presentations/agu_talk_H138.pptx  (gitignored)
"""
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

REPO = Path(__file__).resolve().parents[1]
FIGS = REPO / "draft" / "overleaf" / "Figures"      # manuscript-quality figures
OUTC = REPO / "figures" / "outcomes"
OUT = REPO / "presentations" / "agu_talk_H138.pptx"

# 16:9
SW, SH = Inches(13.333), Inches(7.5)

INK = RGBColor(0x1F, 0x2D, 0x3A)   # near-black headline
MUTE = RGBColor(0x6B, 0x7A, 0x88)  # muted subtitle
ACC = RGBColor(0x2E, 0x68, 0x99)   # project steel blue
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def _txt(slide, text, left, top, width, height, size, color, bold=False,
         align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return box


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def notes(slide, text):
    """Attach speaker notes to a slide."""
    slide.notes_slide.notes_text_frame.text = text
    return slide


def title_slide(prs, title, subtitle, byline):
    s = blank(prs)
    _txt(s, title, Inches(1.0), Inches(2.15), SW - Inches(2.0), Inches(2.0),
         40, INK, bold=True)
    _txt(s, subtitle, Inches(1.0), Inches(4.15), SW - Inches(2.0), Inches(0.7),
         20, ACC)
    _txt(s, byline, Inches(1.0), Inches(5.15), SW - Inches(2.0), Inches(1.0),
         14, MUTE)
    return s


def statement_slide(prs, headline, sub=None):
    """A near-empty slide: one line that lands."""
    s = blank(prs)
    _txt(s, headline, Inches(1.0), Inches(2.6), SW - Inches(2.0), Inches(2.0),
         34, INK, bold=True)
    if sub:
        _txt(s, sub, Inches(1.0), Inches(4.5), SW - Inches(2.0), Inches(1.0),
             18, MUTE)
    return s


def stat_slide(prs, stats, headline=None):
    """Big numbers, minimal words. stats = [(number, label), ...]"""
    s = blank(prs)
    if headline:
        _txt(s, headline, Inches(0.9), Inches(0.8), SW - Inches(1.8), Inches(0.8),
             28, INK, bold=True)
    n = len(stats)
    colw = (SW - Inches(1.8)) / n
    for i, (num, lab) in enumerate(stats):
        left = Inches(0.9) + int(colw * i)
        _txt(s, num, left, Inches(2.5), int(colw), Inches(1.4), 60, ACC,
             bold=True, align=PP_ALIGN.CENTER)
        _txt(s, lab, left, Inches(4.0), int(colw), Inches(1.2), 15, MUTE,
             align=PP_ALIGN.CENTER)
    return s


def figure_slide(prs, img, headline, sub=None):
    """Takeaway as the title; figure fills the rest. Headline and subtitle
    heights adapt to their text length, so a title that wraps to two lines
    no longer collides with the caption or the figure."""
    s = blank(prs)
    text_w = SW - Inches(1.4)

    # Headline (27 pt bold): allow wrapping to multiple lines (~60 chars/line).
    hl_lines = max(1, -(-len(headline) // 60))
    hl_h = 0.52 * hl_lines + 0.10
    _txt(s, headline, Inches(0.7), Inches(0.42), text_w, Inches(hl_h),
         27, INK, bold=True)
    cur = 0.42 + hl_h + 0.06

    if sub:
        # Subtitle (14 pt): ~115 chars/line at this width.
        sub_lines = max(1, -(-len(sub) // 115))
        sub_h = 0.26 * sub_lines + 0.06
        _txt(s, sub, Inches(0.7), Inches(cur), text_w, Inches(sub_h), 14, MUTE)
        cur += sub_h + 0.10

    top = Inches(cur)
    max_w = text_w
    max_h = SH - top - Inches(0.45)
    iw, ih = Image.open(img).size
    scale = min(max_w / iw, max_h / ih)
    w, h = int(iw * scale), int(ih * scale)
    s.shapes.add_picture(str(img), int((SW - w) / 2), top, width=w, height=h)
    return s


def main():
    prs = Presentation()
    prs.slide_width, prs.slide_height = SW, SH

    # 1 — Title
    title_slide(
        prs,
        "Legacy earthworks are an unmeasured human control on the dryland water cycle",
        "775 water spreader berms, Altar Valley, Arizona — quantified with Sentinel-2",
        "O.V. Crompton, A.N. Koop, S. Anderson, S. Chaulagain, D.A. Lapides, "
        "M.B. Meles, M.H. Nichols\nUSDA-ARS   |   AGU H138 — Understanding and "
        "Quantifying the Human Impacts on the Water Cycle",
    )

    # 2 — Hook: the human alteration
    #     (former slide 3 "the gap" now lives in these speaker notes)
    notes(figure_slide(
        prs, FIGS / "LocMap.png",                      # manuscript Figure 1
        "In the 1930s we re-engineered these watersheds — then walked away",
        "More than 1,500 structures built across the 247,000 ha Altar Valley since "
        "the early 1900s; we analyze 775 mapped water spreader berms "
        "(green, intact; red, degraded)",
    ),
        "Thousands of these structures still redirect water. Almost none are "
        "represented in our models. Their condition and their hydrologic effect "
        "are essentially unmeasured at scale.")

    # 3 — Method
    #     (former slide 5 "scale of the effect" now lives in these speaker notes)
    notes(figure_slide(
        prs, FIGS / "fig2_veg_response.png",
        "Satellite greenness turns each berm into a measurement",
        "ΔS = percent SAVI difference 15–60 m upslope vs. downslope, normalized by "
        "background >200 m from any berm (Aug–Sep 2016–2024, Sentinel-2)",
    ),
        "Ninety years on, the outcomes are mixed: 775 berms analyzed; 41% "
        "structurally degraded (203 flanked, 114 breached); 47% show a vegetation "
        "response (ΔS > 7%).")

    # 4 — Condition controls
    figure_slide(
        prs, FIGS / "fig6_berm_condition.png",
        "Structural condition is set by length, flow accumulation, and soil texture",
        "Shorter berms, in positions of lower flow accumulation, on finer-textured "
        "soils stay intact (ΔAIC +30.6, +6.7, +3.0)",
    )

    # 5 — Vegetation response controls
    figure_slide(
        prs, FIGS / "fig7_vegetation_response.png",
        "Vegetation response is set by slope and soil texture — a different recipe",
        "Steeper slopes and coarser-textured (sandy loam) soils give the largest "
        "response (ΔAIC +18.8 and +5.6; RF CV AUC 0.71 condition / 0.64 response)",
    )

    # 6 — Key finding
    figure_slide(
        prs, FIGS / "fig8_veg_response_by_condition_texture.png",
        "An intact berm is not a working berm",
        "Structural condition and vegetation response are statistically "
        "independent (φ = −0.02, p = 0.612); only sandy-loam berms differ "
        "(52% vs. 38%, p = 0.034)",
    )

    # 7 — Confounding / controlled analysis
    #     (former slide 10 "so what" now lives in these speaker notes)
    notes(figure_slide(
        prs, FIGS / "fig3_controlled_predictors.png",
        "Landform looks important — until you control for slope and soil",
        "Univariately associated with vegetation response (p < 0.001), landform "
        "drops to p = 0.310 once slope and soil texture are controlled",
    ),
        "To represent human influence, it is not enough to know that a structure "
        "exists — or even that it is intact. Function depends on landscape context, "
        "and remote sensing can measure it at scale.")

    # 8 — Takeaways
    s = blank(prs)
    _txt(s, "Takeaways", Inches(0.9), Inches(0.9), SW - Inches(1.8), Inches(0.8),
         30, INK, bold=True)
    points = [
        "Legacy earthworks are a widespread, unrepresented human control on dryland hydrology.",
        "Structural condition and vegetation response are shaped by largely separate landscape controls.",
        "Condition is determined by berm length, flow accumulation, and soil texture; vegetation response by slope and soil texture.",
        "The two outcomes are statistically independent — physical integrity and ecological function must be evaluated separately.",
        "These controls can be read from existing terrain and soil-survey data, guiding maintenance or removal at scale.",
    ]
    y = Inches(2.1)
    for p in points:
        _txt(s, "—", Inches(1.0), y, Inches(0.4), Inches(0.6), 18, ACC, bold=True)
        _txt(s, p, Inches(1.5), y, SW - Inches(2.6), Inches(0.9), 18, INK)
        y += Inches(1.0)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"Saved {OUT}  ({len(prs.slides._sldIdLst)} slides)")


if __name__ == "__main__":
    main()
