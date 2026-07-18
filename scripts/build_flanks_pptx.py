"""Compile the figures from the two flanks/breaches analysis notebooks into a
single PowerPoint deck.

Figures are read from ``figures/failure_mechanisms/`` and captions from that
folder's ``figure_registry_concise.txt``. One figure per slide, grouped by the
notebook that produced it.

Usage:
    python scripts/build_flanks_pptx.py
Output:
    presentations/flanks_breaches_figures.pptx  (gitignored)
"""
from pathlib import Path
import re

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

REPO = Path(__file__).resolve().parents[1]
FIG_DIR = REPO / "figures" / "failure_mechanisms"
OUT = REPO / "presentations" / "flanks_breaches_figures.pptx"

# Figures in the order each notebook produces them.
MAIN_NB = "analysis - flanks breaches.ipynb"
PORT_NB = "analysis - flanks breaches port.ipynb"

MAIN_FIGS = [
    "fig3_texture_fail_type.png",
    "fig4_landform_fail_type.png",
    "fig5_length_fail_type.png",
    "fig6_slope_fail_type.png",
    "fig7_soildev_fail_type.png",
    "fig_stacked_fail_by_landform.png",
    "fig_angle_fail_type.png",
    "fig_flowaccum_fail_type.png",
    "fig_savi_fail_type.png",
    "fig_roaddist_fail_type.png",
    "fig13_highclay_fail_type.png",
]
PORT_FIGS = [
    "fig2_length_fail_type_brackets.png",
    "fig2b_slope_fail_type_brackets.png",
    "fig2c_soildev_fail_type_brackets.png",
    "fig2d_texture_fail_type.png",
    "fig3_landform_fail_type.png",
]

# Analysis topic per figure — used to drop replicate figures that appear in
# both notebooks (the port reproduces the 5 univariate panels). First
# occurrence wins, so the comprehensive notebook's versions are kept.
TOPIC = {
    "fig3_texture_fail_type.png": "texture",
    "fig4_landform_fail_type.png": "landform",
    "fig5_length_fail_type.png": "length",
    "fig6_slope_fail_type.png": "slope",
    "fig7_soildev_fail_type.png": "soildev",
    "fig_stacked_fail_by_landform.png": "landform_composition",
    "fig_angle_fail_type.png": "angle",
    "fig_flowaccum_fail_type.png": "flowaccum",
    "fig_savi_fail_type.png": "savi",
    "fig_roaddist_fail_type.png": "roaddist",
    "fig13_highclay_fail_type.png": "highclay",
    "fig2_length_fail_type_brackets.png": "length",
    "fig2b_slope_fail_type_brackets.png": "slope",
    "fig2c_soildev_fail_type_brackets.png": "soildev",
    "fig2d_texture_fail_type.png": "texture",
    "fig3_landform_fail_type.png": "landform",
}

# Clean slide titles keyed by filename.
TITLES = {
    "fig3_texture_fail_type.png": "Soil texture × Failure type",
    "fig4_landform_fail_type.png": "Landform × Failure type",
    "fig5_length_fail_type.png": "Berm length × Failure type",
    "fig6_slope_fail_type.png": "Slope × Failure type",
    "fig7_soildev_fail_type.png": "Soil development (B horizon) × Failure type",
    "fig_stacked_fail_by_landform.png": "Failure-type composition by landform",
    "fig_angle_fail_type.png": "Berm angle (directionality) × Failure type",
    "fig_flowaccum_fail_type.png": "Flow accumulation × Failure type",
    "fig_savi_fail_type.png": "SAVI vegetation indices × Failure type",
    "fig_roaddist_fail_type.png": "Road distance × Failure type",
    "fig13_highclay_fail_type.png": "Clay content × Failure type",
    "fig2_length_fail_type_brackets.png": "Berm length × Failure type",
    "fig2b_slope_fail_type_brackets.png": "Slope × Failure type",
    "fig2c_soildev_fail_type_brackets.png": "Soil development (B horizon) × Failure type",
    "fig2d_texture_fail_type.png": "Soil texture × Failure type",
    "fig3_landform_fail_type.png": "Landform × Failure type",
}

# Slide geometry (16:9).
SW, SH = Inches(13.333), Inches(7.5)
NAVY = RGBColor(0x2E, 0x68, 0x99)
GREY = RGBColor(0x55, 0x55, 0x55)


def load_captions():
    """Map figure filename -> interpretation text from the concise registry."""
    reg = FIG_DIR / "figure_registry_concise.txt"
    caps = {}
    if not reg.exists():
        return caps
    fname = None
    for line in reg.read_text().splitlines():
        m = re.match(r"\s*File\s*:\s*(\S+)", line)
        if m:
            fname = m.group(1)
        m = re.match(r"\s*Interpretation:\s*(.+)", line)
        if m and fname:
            caps[fname] = m.group(1).strip()
    return caps


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.6), SW - Inches(1.6), Inches(2))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = title
    r.font.size = Pt(40); r.font.bold = True; r.font.color.rgb = NAVY
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = subtitle
    r2.font.size = Pt(20); r2.font.color.rgb = GREY
    return slide


def add_section_slide(prs, text):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(0.8), Inches(3.1), SW - Inches(1.6), Inches(1.3))
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(28); r.font.bold = True; r.font.color.rgb = NAVY
    return slide


def add_figure_slide(prs, img, title, caption):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Title
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), SW - Inches(1.0), Inches(0.8))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title
    r.font.size = Pt(26); r.font.bold = True; r.font.color.rgb = NAVY
    # Image — fit within a box, preserve aspect ratio
    from PIL import Image
    iw, ih = Image.open(img).size
    max_w, max_h = SW - Inches(1.2), Inches(5.2)
    scale = min(max_w / iw, max_h / ih)
    w, h = int(iw * scale), int(ih * scale)
    left = int((SW - w) / 2)
    slide.shapes.add_picture(str(img), left, Inches(1.15), width=w, height=h)
    # Caption
    if caption:
        cb = slide.shapes.add_textbox(Inches(0.6), Inches(6.55), SW - Inches(1.2), Inches(0.85))
        ctf = cb.text_frame; ctf.word_wrap = True
        cp = ctf.paragraphs[0]; cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run(); cr.text = caption
        cr.font.size = Pt(12); cr.font.italic = True; cr.font.color.rgb = GREY
    return slide


def main():
    caps = load_captions()
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH

    add_title_slide(
        prs,
        "Altar Valley Berms — Breach vs Flank Failure Mechanisms",
        "Figures compiled from the two flanks/breaches analysis notebooks",
    )

    n_added = 0
    seen_topics = set()
    for nb, figs, label in [
        (MAIN_NB, MAIN_FIGS, "Comprehensive analysis"),
        (PORT_NB, PORT_FIGS, "Bracketed univariate panels (port notebook)"),
    ]:
        # Figures for this section that aren't replicates of one already added.
        section = []
        for fn in figs:
            topic = TOPIC.get(fn, fn)
            if topic in seen_topics:
                print(f"  [dedup] skip replicate: {fn} (topic '{topic}')")
                continue
            if not (FIG_DIR / fn).exists():
                print(f"  [skip] missing: {fn}")
                continue
            section.append(fn)
            seen_topics.add(topic)
        if not section:
            print(f"  [dedup] section '{label}' empty after dedup — omitted")
            continue
        add_section_slide(prs, f"{label}\n{nb}")
        for fn in section:
            add_figure_slide(prs, FIG_DIR / fn, TITLES.get(fn, fn), caps.get(fn, ""))
            n_added += 1

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"Saved {OUT}  ({n_added} figure slides, {len(prs.slides.__iter__.__self__._sldIdLst)} total slides)")


if __name__ == "__main__":
    main()
